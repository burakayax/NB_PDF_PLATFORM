"""
Ek PDF araçları (sayfa yönetimi, dönüşüm, güvenlik).
pymupdf (fitz), pikepdf, PyPDF2 ve mevcut pdf_engine yardımcılarıyla.
"""

from __future__ import annotations

import io
import os
import shutil
import tempfile
import zipfile
from typing import Dict, List, Optional

import fitz

# pdf_engine'den parola açma
from src.pdf_engine import _open_pdf_reader, is_pdf_encrypted, get_num_pages

# Web SaaS: single quality tier (DPI not user-configurable).
PDF_EXPORT_DPI_WEB = 300
_RASTER_PAGE_BATCH = 6


def _fitz_open(pdf_path: str, password: Optional[str] = None):
    doc = fitz.open(pdf_path)
    if doc.needs_pass:
        if not (password or "").strip():
            raise Exception("Şifreli PDF için parola gerekli.")
        if not doc.authenticate(password or ""):
            raise Exception("Girilen PDF parolası hatalı.")
    return doc


def delete_pages_pdf(
    pdf_path: str,
    output_path: str,
    pages_to_delete: List[int],
    password: Optional[str] = None,
    *,
    total_pages: Optional[int] = None,
) -> bool:
    """İstenen sayfalar çıkarılmış yeni belge oluşturur (insert_pdf aralık aktarımı ile büyük PDF’lerde çok daha hızlı)."""
    import fitz as _fitz_local
    src = _fitz_open(pdf_path, password=password)
    try:
        n = src.page_count
        to_del = {int(p) for p in pages_to_delete}
        if any(p < 1 or p > n for p in to_del):
            raise Exception("Geçersiz sayfa numarası.")
        if len(to_del) >= n:
            raise Exception("Tüm sayfalar silinemez; en az bir sayfa kalmalıdır.")
        keep = [i for i in range(n) if (i + 1) not in to_del]
        new_doc = _fitz_local.open()
        try:
            ranges: list[tuple[int, int]] = []
            if keep:
                s = keep[0]; e = keep[0]
                for k in keep[1:]:
                    if k == e + 1:
                        e = k
                    else:
                        ranges.append((s, e))
                        s = e = k
                ranges.append((s, e))
            for from_p, to_p in ranges:
                new_doc.insert_pdf(src, from_page=from_p, to_page=to_p)
            new_doc.save(output_path, garbage=0, deflate=False, linear=False)
        finally:
            new_doc.close()
    finally:
        src.close()
    return True


def rotate_pdf(
    pdf_path: str,
    output_path: str,
    degrees: int,
    pages_1based: Optional[List[int]],
    password: Optional[str] = None,
    per_page_degrees: Optional[Dict[int, int]] = None,
) -> bool:
    """Rotate pages. Use either legacy (degrees + pages_1based) or per_page_degrees (1-based page -> add 90/180/270)."""
    doc = _fitz_open(pdf_path, password=password)
    try:
        n = doc.page_count
        if per_page_degrees is not None:
            for p in range(1, n + 1):
                add_deg = int(per_page_degrees.get(p, 0))
                if add_deg == 0:
                    continue
                if add_deg not in (90, 180, 270):
                    raise Exception("Sayfa başına dönüş yalnızca 90, 180 veya 270 olabilir.")
                page = doc[p - 1]
                cur = int(page.rotation) % 360
                page.set_rotation((cur + add_deg) % 360)
        else:
            if degrees not in (90, 180, 270):
                raise Exception("Dönüş açısı 90, 180 veya 270 olmalıdır.")
            targets = [p - 1 for p in (pages_1based or list(range(1, n + 1)))]
            for i in targets:
                if i < 0 or i >= n:
                    continue
                page = doc[i]
                cur = int(page.rotation) % 360
                page.set_rotation((cur + degrees) % 360)
        # Döndürme yalnızca sayfa sözlüğündeki /Rotate meta-verisini değiştirir;
        # içerik akışları dokunulmaz — yeniden sıkıştırmaya gerek yok.
        doc.save(output_path, garbage=0, deflate=False, linear=False)
    finally:
        doc.close()
    return True


def organize_pdf(
    pdf_path: str,
    output_path: str,
    new_order_1based: List[int],
    password: Optional[str] = None,
    *,
    total_pages: Optional[int] = None,
) -> bool:
    """Sayfaları yeni sıraya göre düzenler (ör. [3,1,2]).

    insert_pdf ile ardışık aralık aktarımı kullanır — büyük PDF'lerde
    select()+save() tam dosya yeniden yazımından çok daha hızlıdır.
    """
    import fitz as _fitz_local
    src = _fitz_open(pdf_path, password=password)
    try:
        n = src.page_count
        for p in new_order_1based:
            if p < 1 or p > n:
                raise Exception(f"Geçersiz sayfa: {p} (1–{n})")
        if len(new_order_1based) != n:
            raise Exception("Sıra listesi, tüm sayfaları tam olarak bir kez içermelidir.")
        if len(set(new_order_1based)) != n:
            raise Exception("Aynı sayfa iki kez kullanılamaz.")
        order_0 = [p - 1 for p in new_order_1based]
        ranges: list[tuple[int, int]] = []
        s = order_0[0]; e = order_0[0]
        for k in order_0[1:]:
            if k == e + 1:
                e = k
            else:
                ranges.append((s, e))
                s = e = k
        ranges.append((s, e))
        new_doc = _fitz_local.open()
        try:
            for from_p, to_p in ranges:
                new_doc.insert_pdf(src, from_page=from_p, to_page=to_p)
            new_doc.save(output_path, garbage=0, deflate=False, linear=False)
        finally:
            new_doc.close()
    finally:
        src.close()
    return True


def unlock_pdf_pikepdf(input_path: str, output_path: str, password: str) -> bool:
    """Kullanıcı parolası ile açıp şifresiz PDF kaydeder."""
    try:
        import pikepdf
    except ImportError as e:
        raise Exception("pikepdf gerekli.") from e
    if not (password or "").strip():
        raise Exception("PDF şifresini girmeniz gerekir.")
    if not is_pdf_encrypted(input_path):
        shutil.copy2(input_path, output_path)
        return True
    with pikepdf.open(input_path, password=password.strip(), allow_overwriting_input=False) as pdf:
        pdf.save(output_path, linearize=True)
    return True


def _hex_to_rgb(hex_color: str) -> tuple:
    """#RRGGBB → (r, g, b) 0-1 aralığında."""
    h = hex_color.lstrip("#")
    if len(h) != 6:
        return (0.55, 0.55, 0.55)
    try:
        return tuple(int(h[i : i + 2], 16) / 255.0 for i in (0, 2, 4))
    except ValueError:
        return (0.55, 0.55, 0.55)


def add_watermark_text(
    input_path: str,
    output_path: str,
    text: str,
    opacity: float = 0.12,
    password: Optional[str] = None,
    font_name: str = "helv",
    font_color: str = "#8C8C8C",
) -> bool:
    if not (text or "").strip():
        raise Exception("Filigran metni boş olamaz.")
    op = max(0.01, min(0.5, float(opacity)))
    color = _hex_to_rgb(font_color)
    valid_fonts = {"helv", "tiro", "cour", "zadb", "symb"}
    fn = font_name if font_name in valid_fonts else "helv"
    doc = _fitz_open(input_path, password=password)
    try:
        for i in range(doc.page_count):
            page = doc[i]
            r = page.rect
            c = fitz.Point((r.x0 + r.x1) / 2, (r.y0 + r.y1) / 2)
            page.insert_text(
                c,
                (text or "").strip(),
                fontname=fn,
                fontsize=22,
                color=color,
                render_mode=0,
                fill_opacity=op,
            )
        doc.save(output_path, garbage=0, deflate=False, linear=False)
    finally:
        doc.close()
    return True


def add_page_numbers(
    input_path: str,
    output_path: str,
    start_at: int = 1,
    position: str = "footer",
    password: Optional[str] = None,
    fmt: str = "plain",
) -> bool:
    """
    fmt: "plain"  → "3"
         "page"   → "Sayfa 3"  / "Page 3"
         "of"     → "3 / 10"
    """
    doc = _fitz_open(input_path, password=password)
    try:
        total = doc.page_count
        num = int(start_at)
        margin_x = 36
        strip_h = 24
        for i in range(total):
            page = doc[i]
            r = page.rect
            if fmt == "page":
                label = f"Sayfa {num}"
            elif fmt == "of":
                label = f"{num} / {total}"
            else:
                label = str(num)
            # Use a full-width rect so insert_textbox can centre the text.
            if position == "header":
                rect = fitz.Rect(r.x0 + margin_x, r.y0 + 6, r.x1 - margin_x, r.y0 + strip_h)
            else:
                rect = fitz.Rect(r.x0 + margin_x, r.y1 - strip_h, r.x1 - margin_x, r.y1 - 6)
            page.insert_textbox(
                rect, label,
                fontsize=9,
                color=(0.4, 0.4, 0.4),
                align=fitz.TEXT_ALIGN_CENTER,
            )
            num += 1
        doc.save(output_path, garbage=0, deflate=False, linear=False)
    finally:
        doc.close()
    return True


def repair_pdf(input_path: str, output_path: str, password: Optional[str] = None) -> bool:
    """Bozuk PDF'i çok aşamalı strateji ile onarır. Tüm yöntemler başarısız olursa açıklayıcı hata verir."""
    try:
        import pikepdf
    except ImportError as e:
        raise Exception("pikepdf kütüphanesi bulunamadı; sunucu yapılandırmasını kontrol edin.") from e

    op = (password or "").strip()
    errors: list[str] = []

    # Strateji 1: pikepdf — çapraz referans tablosunu yeniden oluşturur, bozuk akışları atlar
    try:
        with pikepdf.open(input_path, password=op, suppress_warnings=True) as pdf:
            pdf.save(output_path, compress_streams=True, recompress_flate=True)
        if os.path.isfile(output_path) and os.path.getsize(output_path) > 32:
            return True
    except pikepdf.PasswordError:
        raise Exception("PDF şifreli; onarım için doğru parolayı girin.")
    except Exception as e1:
        errors.append(f"pikepdf: {e1!s:.150}")

    # Strateji 2: PyMuPDF — xref tablosunu yeniden oluşturur, stream'leri temizler
    try:
        doc = _fitz_open(input_path, password=password)
        try:
            doc.save(output_path, garbage=4, deflate=True, clean=True, linear=False)
        finally:
            doc.close()
        if os.path.isfile(output_path) and os.path.getsize(output_path) > 32:
            return True
    except Exception as e2:
        if "password" in str(e2).lower() or "encrypted" in str(e2).lower():
            raise Exception("PDF şifreli; onarım için doğru parolayı girin.")
        errors.append(f"fitz: {e2!s:.150}")

    # Strateji 3: pikepdf lenient mode (çok bozuk dosyalar için, kurtarılabilir sayfalar)
    try:
        with pikepdf.open(input_path, password=op, suppress_warnings=True, ignore_xref_streams=True) as pdf:
            if len(pdf.pages) == 0:
                errors.append("pikepdf-lenient: sayfa bulunamadı")
            else:
                pdf.save(output_path, compress_streams=True)
                if os.path.isfile(output_path) and os.path.getsize(output_path) > 32:
                    return True
    except Exception as e3:
        errors.append(f"pikepdf-lenient: {e3!s:.150}")

    err_summary = " | ".join(errors[:3])
    raise Exception(
        f"PDF onarılamadı. Dosya kurtarılamayacak kadar ciddi biçimde bozulmuş olabilir. "
        f"Orijinal dosyanın yedeği varsa onu kullanın. ({err_summary})"
    )


def pdf_to_text(input_path: str, output_path: str, password: Optional[str] = None) -> bool:
    """PDF içindeki metin katmanını düz metin dosyasına yazar (sayfa başlıkları dahil)."""
    doc = _fitz_open(input_path, password=password)
    try:
        lines: list[str] = []
        for page_num, page in enumerate(doc, 1):
            text = page.get_text("text").strip()
            if text:
                lines.append(f"--- Sayfa {page_num} ---")
                lines.append(text)
        if not lines:
            raise Exception("PDF içinde çıkarılabilir metin bulunamadı. Taranmış görüntü PDF'leri metin içermez.")
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n".join(lines))
    finally:
        doc.close()
    return True


def flatten_pdf(input_path: str, output_path: str, password: Optional[str] = None) -> bool:
    """Etkileşimli form alanlarını ve açıklamaları PDF içeriğine gömer (düzleştirir)."""
    try:
        import pikepdf
    except ImportError as e:
        raise Exception("pikepdf kütüphanesi bulunamadı.") from e
    op = (password or "").strip()
    try:
        # pikepdf password varsayılanı "" (parolasız); None geçmek bazı sürümlerde
        # TypeError verir. Her zaman string geçiyoruz.
        with pikepdf.open(input_path, password=op) as pdf:
            # Form alanlarını sil — içerik zaten sayfaya render edilmiş olarak kalır
            if "/AcroForm" in pdf.Root:
                del pdf.Root["/AcroForm"]
            for page in pdf.pages:
                if "/Annots" in page:
                    annots = page["/Annots"]
                    keep = []
                    for annot in annots:
                        atype = str(annot.get("/Subtype", ""))
                        if atype not in ("/Widget", "/FreeText", "/Stamp", "/Highlight",
                                         "/Underline", "/StrikeOut", "/Squiggly", "/Caret"):
                            keep.append(annot)
                    if keep:
                        page["/Annots"] = pikepdf.Array(keep)
                    else:
                        del page["/Annots"]
            pdf.save(output_path, compress_streams=True)
    except pikepdf.PasswordError:
        raise Exception("PDF şifreli; düzleştirmek için doğru parolayı girin.")
    except Exception as e:
        if "password" in str(e).lower():
            raise Exception("PDF şifreli; düzleştirmek için doğru parolayı girin.")
        raise
    return True


def pdf_to_images_zip(
    pdf_path: str,
    workdir: str,
    image_format: str = "jpg",
    dpi: int = PDF_EXPORT_DPI_WEB,
    password: Optional[str] = None,
) -> str:
    """ZIP dosya yolunu döndürür; pdf2image + Poppler gerekir. Sayfalar partiler halinde rasterize edilir (bellek)."""
    from pdf2image import convert_from_path

    fmt = (image_format or "jpg").lower()
    if fmt not in ("jpg", "jpeg", "png"):
        raise Exception("Görüntü formatı jpg veya png olmalıdır.")
    ext = "png" if fmt == "png" else "jpg"
    import src.pdf_engine as pe

    import os as _os
    _cpu = max(1, min((_os.cpu_count() or 2), 4))
    poppler = getattr(pe, "poppler_bin_path", None) or None
    kw_base: dict = {"dpi": int(dpi), "fmt": "png" if ext == "png" else "jpeg", "thread_count": _cpu}
    if poppler and os.path.isdir(poppler):
        kw_base["poppler_path"] = poppler
    pwd = (password or "").strip()
    if pwd:
        kw_base["userpw"] = pwd
    _open_pdf_reader(pdf_path, password=password)
    n = get_num_pages(pdf_path, password=password)
    zip_path = os.path.join(workdir, "sayfalar.zip")
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
        page_index = 0
        for start in range(1, n + 1, _RASTER_PAGE_BATCH):
            end = min(start + _RASTER_PAGE_BATCH - 1, n)
            kw = {**kw_base, "first_page": start, "last_page": end}
            images = convert_from_path(pdf_path, **kw)
            for im in images:
                page_index += 1
                buf = io.BytesIO()
                if ext == "png":
                    im.save(buf, format="PNG")
                    name = f"page_{page_index:04d}.png"
                else:
                    im.save(buf, format="JPEG", quality=90)
                    name = f"page_{page_index:04d}.jpg"
                zf.writestr(name, buf.getvalue())
                del im
    return zip_path


def images_to_pdf(image_paths: List[str], output_path: str) -> bool:
    try:
        import img2pdf
    except ImportError:
        merged = fitz.open()
        try:
            for p in image_paths:
                imgdoc = fitz.open(p)
                try:
                    pdfb = imgdoc.convert_to_pdf()
                finally:
                    imgdoc.close()
                m = fitz.open("pdf", pdfb)
                try:
                    merged.insert_pdf(m)
                finally:
                    m.close()
            merged.save(output_path, garbage=1, deflate=False)
        finally:
            merged.close()
        return True
    with open(output_path, "wb") as f:
        f.write(img2pdf.convert(image_paths))
    return True


def html_to_pdf_file(html: str, output_path: str, base_url: Optional[str] = None) -> bool:
    """Önce wkhtmltopdf (daha uyumlu), sonra xhtml2pdf dener.

    wkhtmltopdf sistemde yoksa xhtml2pdf ile devam eder (sınırlı CSS desteği).
    Her ikisi de başarısız olursa kullanıcı dostu bir hata mesajı fırlatır.
    """
    wk = shutil.which("wkhtmltopdf")
    if wk:
        try:
            import pdfkit

            cfg = pdfkit.configuration(wkhtmltopdf=wk)
            pdfkit.from_string(
                html or "<html><body></body></html>",
                output_path,
                configuration=cfg,
                options={
                    "quiet": "",
                    "enable-local-file-access": "",
                    "disable-smart-shrinking": "",
                },
            )
            if os.path.isfile(output_path) and os.path.getsize(output_path) > 32:
                return True
        except Exception:
            if os.path.isfile(output_path):
                try:
                    os.remove(output_path)
                except OSError:
                    pass
    try:
        from xhtml2pdf import pisa
    except ImportError as e:
        raise Exception(
            "HTML→PDF dönüşümü şu anda kullanılamıyor. "
            "Sunucuda wkhtmltopdf kurulu değil ve xhtml2pdf paketi de bulunamadı. "
            "Lütfen daha sonra tekrar deneyin veya destek ekibiyle iletişime geçin."
        ) from e
    html_src = html or "<html><body></body></html>"
    # xhtml2pdf modern CSS'i parse edemez; tüm stylesheet referanslarını ve inline style bloklarını kaldır
    import re as _re
    html_src = _re.sub(r'<link[^>]+rel=["\']stylesheet["\'][^>]*>', '', html_src, flags=_re.IGNORECASE)
    html_src = _re.sub(r'<link[^>]+href=["\'][^"\']*\.css[^"\']*["\'][^>]*>', '', html_src, flags=_re.IGNORECASE)
    html_src = _re.sub(r'<style[^>]*>.*?</style>', '', html_src, flags=_re.IGNORECASE | _re.DOTALL)
    html_src = _re.sub(r'@import\s+["\'][^"\']+["\'];?', '', html_src, flags=_re.IGNORECASE)
    if "<meta charset" not in html_src[:1000].lower():
        if "<head>" in html_src.lower():
            html_src = html_src.replace("<head>", '<head><meta charset="utf-8">', 1)
        else:
            html_src = f'<html><head><meta charset="utf-8"></head><body>{html_src}</body></html>'
    with open(output_path, "wb") as out:
        status = pisa.CreatePDF(
            src=html_src.encode("utf-8"),
            dest=out,
            encoding="utf-8",
            path_base=base_url or None,
        )
    if status.err or not (os.path.isfile(output_path) and os.path.getsize(output_path) > 32):
        raise Exception(
            "HTML PDF'e dönüştürülemedi. "
            "Sayfanın geçerli HTML içerdiğinden emin olun ve JavaScript gerektirmeyen basit sayfalar deneyin. "
            "Karmaşık CSS/JS içeren sayfalar için URL yerine sayfa kaynağını (HTML metnini) kullanın."
        )
    return True


def html_url_to_pdf(url: str, output_path: str) -> bool:
    import httpx

    u = (url or "").strip()
    if not u.startswith(("http://", "https://")):
        u = "https://" + u
    r = httpx.get(u, timeout=60.0, follow_redirects=True)
    r.raise_for_status()
    ct = r.headers.get("content-type", "")
    if "html" not in ct.lower() and "text" not in ct.lower() and "application" not in ct.lower():
        pass
    return html_to_pdf_file(r.text, output_path, base_url=u)


def pdf_to_pptx(pdf_path: str, pptx_path: str, password: Optional[str] = None, dpi: int = PDF_EXPORT_DPI_WEB) -> bool:
    from pdf2image import convert_from_path
    from pptx import Presentation
    from pptx.util import Emu
    import pdfplumber
    import src.pdf_engine as pe

    import os as _os2
    _cpu2 = max(1, min((_os2.cpu_count() or 2), 4))
    poppler = getattr(pe, "poppler_bin_path", None) or None
    kw_base: dict = {"dpi": int(dpi), "fmt": "png", "thread_count": _cpu2}
    if poppler and os.path.isdir(poppler):
        kw_base["poppler_path"] = poppler
    pwd = (password or "").strip()
    if pwd:
        kw_base["userpw"] = pwd
    _open_pdf_reader(pdf_path, password=password)
    n = get_num_pages(pdf_path, password=password)

    # PDF'in gerçek sayfa boyutunu al (ilk sayfa referans)
    # pdfplumber sayfa boyutunu pt cinsinden verir; 1 pt = 12700 EMU
    PT_TO_EMU = 12700
    page_w_pt, page_h_pt = 595.0, 842.0  # A4 varsayılan
    try:
        with pdfplumber.open(pdf_path, password=pwd or "") as _pdf:
            if _pdf.pages:
                p0 = _pdf.pages[0]
                page_w_pt = float(p0.width)
                page_h_pt = float(p0.height)
    except Exception:
        pass

    prs = Presentation()
    # Slayt boyutunu PDF sayfa boyutuna eşitle → görüntü bozulmadan yerleşir
    prs.slide_width = Emu(int(page_w_pt * PT_TO_EMU))
    prs.slide_height = Emu(int(page_h_pt * PT_TO_EMU))

    try:
        blank = prs.slide_layouts[6]
    except (IndexError, KeyError):
        blank = prs.slide_layouts[0]

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    for start in range(1, n + 1, _RASTER_PAGE_BATCH):
        end = min(start + _RASTER_PAGE_BATCH - 1, n)
        kw = {**kw_base, "first_page": start, "last_page": end}
        images = convert_from_path(pdf_path, **kw)
        for im in images:
            slide = prs.slides.add_slide(blank)
            fd, tmp = tempfile.mkstemp(suffix=".png")
            os.close(fd)
            try:
                im.save(tmp, "PNG")
                img_w_px, img_h_px = im.size
                # En-boy oranını koruyarak slayta sığdır (letterbox)
                scale = min(slide_w / img_w_px, slide_h / img_h_px)
                pic_w = Emu(int(img_w_px * scale))
                pic_h = Emu(int(img_h_px * scale))
                left = Emu(int((slide_w - pic_w) / 2))
                top = Emu(int((slide_h - pic_h) / 2))
                slide.shapes.add_picture(tmp, left, top, width=pic_w, height=pic_h)
            finally:
                try:
                    os.remove(tmp)
                except OSError:
                    pass
            del im
    prs.save(pptx_path)
    return True


def pptx_to_pdf(pptx_path: str, pdf_path: str) -> bool:
    import shutil
    import subprocess

    timeout_sec = max(30, int(os.environ.get("NB_PDF_TOOL_TIMEOUT_SEC", "300")))

    def _soffice_executable() -> Optional[str]:
        for c in (shutil.which("soffice"), shutil.which("libreoffice")):
            if c:
                return c
        if os.name == "nt":
            for pf in (
                os.environ.get("ProgramFiles", r"C:\Program Files"),
                os.environ.get("ProgramFiles(x86)", ""),
            ):
                if not pf:
                    continue
                for sub in (
                    os.path.join(pf, "LibreOffice", "program", "soffice.com"),
                    os.path.join(pf, "LibreOffice", "program", "soffice.exe"),
                ):
                    if os.path.isfile(sub):
                        return sub
        return None

    def _via_libreoffice() -> bool:
        soffice = _soffice_executable()
        if not soffice:
            return False
        pptx_abs = os.path.abspath(pptx_path)
        out_dir = os.path.dirname(os.path.abspath(pdf_path))
        os.makedirs(out_dir, exist_ok=True)
        try:
            subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf", "--outdir", out_dir, pptx_abs],
                check=True,
                timeout=timeout_sec,
                capture_output=True,
                text=True,
            )
        except subprocess.CalledProcessError as e:
            raise Exception(f"LibreOffice dönüşümü başarısız: {e.stderr or e}") from e
        base = os.path.splitext(os.path.basename(pptx_abs))[0]
        produced = os.path.join(out_dir, base + ".pdf")
        if not os.path.isfile(produced):
            raise Exception("LibreOffice çıktı dosyası oluşmadı.")
        if os.path.abspath(produced) != os.path.abspath(pdf_path):
            shutil.move(produced, pdf_path)
        return os.path.isfile(pdf_path)

    if os.name == "nt":
        try:
            import pythoncom
            from win32com.client import DispatchEx

            com_inited = False
            try:
                pythoncom.CoInitialize()
                com_inited = True
            except Exception:
                com_inited = False
            app = None
            pres = None
            try:
                app = DispatchEx("PowerPoint.Application")
                try:
                    app.DisplayAlerts = 0
                except Exception:
                    pass
                pres = app.Presentations.Open(os.path.abspath(pptx_path), WithWindow=False, ReadOnly=True)
                out = os.path.abspath(pdf_path)
                pres.SaveAs(out, 32)  # ppSaveAsPDF
                pres.Close()
                pres = None
            finally:
                if pres is not None:
                    try:
                        pres.Close()
                    except Exception:
                        pass
                if app is not None:
                    try:
                        app.Quit()
                    except Exception:
                        pass
                if com_inited:
                    try:
                        pythoncom.CoUninitialize()
                    except Exception:
                        pass
            return os.path.isfile(pdf_path)
        except Exception:
            if _soffice_executable():
                return _via_libreoffice()
            raise

    if not _soffice_executable():
        raise Exception(
            "PPTX→PDF: LibreOffice gerekli (`soffice` veya `libreoffice` PATH'te). "
            "Windows'ta PowerPoint yüklüyse o da denenir."
        )
    return _via_libreoffice()
